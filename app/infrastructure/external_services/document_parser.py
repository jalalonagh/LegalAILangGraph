"""
Document parser abstraction and universal implementation.

Supports PDF, DOCX, TXT, HTML, and other formats. PDF documents are
parsed with pdfplumber first; if no text is found, OCR is attempted
for scanned documents.
"""

from __future__ import annotations

import asyncio
import html
import io
import os
import re
import shutil
import tempfile
from typing import Any

from app.core.config import settings
from app.core.logging import get_logger
from app.core.exceptions import UnsupportedFileTypeError
from app.domain.enums import DocumentParseFormat
from app.domain.interfaces import DocumentParser, OCRProvider, ParsedDocument
from app.domain.value_objects import classify_file_format

logger = get_logger()


class UniversalDocumentParser(DocumentParser):
    """Universal document parser supporting multiple formats."""

    def __init__(self, ocr_provider: OCRProvider | None = None) -> None:
        self._ocr = ocr_provider

    @property
    def supported_formats(self) -> list[DocumentParseFormat]:
        return [
            DocumentParseFormat.PDF,
            DocumentParseFormat.DOCX,
            DocumentParseFormat.DOC,
            DocumentParseFormat.TXT,
            DocumentParseFormat.HTML,
            DocumentParseFormat.RTF,
            DocumentParseFormat.XLSX,
            DocumentParseFormat.XLS,
            DocumentParseFormat.PPTX,
            DocumentParseFormat.CSV,
        ]

    async def parse(
        self,
        file_path: str,
        mime_type: str,
        fmt: DocumentParseFormat,
        **kwargs: Any,
    ) -> ParsedDocument:
        """Parse a document file and return text content + metadata."""
        if fmt not in self.supported_formats:
            raise UnsupportedFileTypeError(
                f"Unsupported file format: {fmt}",
                details={"mime_type": mime_type, "format": fmt.value},
            )

        if fmt == DocumentParseFormat.PDF:
            return await self._parse_pdf(file_path)
        elif fmt == DocumentParseFormat.DOCX:
            return await self._parse_docx(file_path)
        elif fmt == DocumentParseFormat.DOC:
            return await self._parse_doc(file_path)
        elif fmt == DocumentParseFormat.TXT:
            return await self._parse_txt(file_path)
        elif fmt == DocumentParseFormat.HTML:
            return await self._parse_html(file_path)
        elif fmt == DocumentParseFormat.RTF:
            return await self._parse_rtf(file_path)
        elif fmt in (DocumentParseFormat.XLSX, DocumentParseFormat.XLS):
            return await self._parse_excel(file_path)
        elif fmt == DocumentParseFormat.PPTX:
            return await self._parse_pptx(file_path)
        elif fmt == DocumentParseFormat.CSV:
            return await self._parse_csv(file_path)
        else:
            raise UnsupportedFileTypeError(f"Unsupported format: {fmt}")

    async def parse_bytes(
        self,
        data: bytes,
        mime_type: str,
        filename: str,
        **kwargs: Any,
    ) -> ParsedDocument:
        """Parse document from raw bytes using a temp file."""
        fmt = classify_file_format(filename)
        if fmt is None:
            raise UnsupportedFileTypeError(
                f"Cannot determine format for {filename}",
                details={"mime_type": mime_type},
            )

        with tempfile.NamedTemporaryFile(
            suffix=f".{fmt.value}", delete=False
        ) as tmp:
            tmp.write(data)
            tmp_path = tmp.name

        try:
            return await self.parse(tmp_path, mime_type, fmt, **kwargs)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    # ------------------------------------------------------------------
    # Format-specific parsers
    # ------------------------------------------------------------------
    async def _parse_pdf(self, file_path: str) -> ParsedDocument:
        text_parts: list[str] = []
        pages: list[str] = []
        tables: list[list[list[str]]] = []

        try:
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                pages_full = pdf.pages
                for i, page in enumerate(pages_full):
                    page_text = page.extract_text() or ""
                    tables_on_page = page.extract_tables() or []
                    text_parts.append(page_text)
                    pages.append(page_text)
                    for table in tables_on_page:
                        tables.append(table)

                page_count = len(pages_full)
                metadata = {
                    "page_count": page_count,
                    "source": os.path.basename(file_path),
                    "parser": "pdfplumber",
                }

                # If pdfplumber found no text, try OCR
                combined = "\n".join(text_parts).strip()
                if not combined and self._ocr and settings.FEATURE_OCR:
                    logger.info("pdf_ocr_attempt", file=file_path)
                    ocr_text = await self._ocr.extract_text(file_path)
                    return ParsedDocument(
                        text=ocr_text,
                        metadata={**metadata, "ocr": True},
                        pages=[ocr_text] if ocr_text else [],
                        tables=tables,
                    )

                return ParsedDocument(text=combined, metadata=metadata, pages=pages or None, tables=tables or None)
        except Exception as exc:
            logger.error("pdf_parse_failed", file=file_path, error=str(exc))
            # Fallback: try OCR
            if self._ocr and settings.FEATURE_OCR:
                ocr_text = await self._ocr.extract_text(file_path)
                return ParsedDocument(
                    text=ocr_text,
                    metadata={"ocr": True, "source": os.path.basename(file_path)},
                    pages=[ocr_text] if ocr_text else [],
                    tables=[],
                )
            return ParsedDocument(text="", metadata={"error": str(exc)}, pages=[], tables=[])

    async def _parse_docx(self, file_path: str) -> ParsedDocument:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        tables = []
        for table in doc.tables:
            tables.append([[cell.text for cell in row.cells] for row in table.rows])
        metadata = {"paragraph_count": len(paragraphs), "source": os.path.basename(file_path)}
        return ParsedDocument(text=text, metadata=metadata, pages=None, tables=tables or None)

    async def _parse_doc(self, file_path: str) -> ParsedDocument:
        # Try LibreOffice conversion to docx, then parse
        libre = shutil.which("libreoffice") or shutil.which("soffice")
        if libre:
            out_dir = tempfile.mkdtemp()
            try:
                cmd = [libre, "--headless", "--convert-to", "docx", "--outdir", out_dir, file_path]
                await _run_subprocess(cmd)
                converted = os.path.join(out_dir, os.path.splitext(os.path.basename(file_path))[0] + ".docx")
                if os.path.exists(converted):
                    return await self._parse_docx(converted)
            except Exception as exc:
                logger.warning("libreoffice_conversion_failed", error=str(exc))
            finally:
                shutil.rmtree(out_dir, ignore_errors=True)

        # Fallback: use antiword if available
        antiword = shutil.which("antiword")
        if antiword:
            try:
                result = await _run_subprocess([antiword, file_path], capture=True)
                text = result.decode("utf-8", errors="replace").strip()
                return ParsedDocument(text=text, metadata={"source": os.path.basename(file_path)}, pages=[], tables=[])
            except Exception as exc:
                logger.warning("antiword_failed", error=str(exc))

        raise UnsupportedFileTypeError(
            "Cannot parse .doc file — no conversion tool available",
            details={"file": os.path.basename(file_path)},
        )

    async def _parse_txt(self, file_path: str) -> ParsedDocument:
        text = await _read_text_file(file_path)
        return ParsedDocument(text=text, metadata={"source": os.path.basename(file_path)}, pages=[], tables=[])

    async def _parse_html(self, file_path: str) -> ParsedDocument:
        text = await _read_text_file(file_path)
        # Strip HTML tags crudely
        clean = re.sub(r"<[^>]+>", " ", text)
        clean = html.unescape(clean)
        clean = re.sub(r"\s+", " ", clean).strip()
        return ParsedDocument(text=clean, metadata={"source": os.path.basename(file_path)}, pages=[], tables=[])

    async def _parse_rtf(self, file_path: str) -> ParsedDocument:
        from striprtf.striprtf import rtf_to_text

        text = await _read_text_file(file_path)
        try:
            clean = rtf_to_text(text)
        except Exception as exc:
            logger.warning("rtf_parse_failed", error=str(exc))
            clean = re.sub(r"<[^>]+>", " ", text)
        return ParsedDocument(text=clean, metadata={"source": os.path.basename(file_path)}, pages=[], tables=[])

    async def _parse_excel(self, file_path: str) -> ParsedDocument:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        texts: list[str] = []
        tables: list[list[list[str]] = []
        for ws in wb.worksheets:
            sheet_rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    texts.append(" | ".join(cells))
                    sheet_rows.append(cells)
            if sheet_rows:
                tables.append(sheet_rows)
        wb.close()
        text = "\n".join(texts)
        return ParsedDocument(
            text=text,
            metadata={"source": os.path.basename(file_path), "sheets": len(tables)},
            pages=None,
            tables=tables,
        )

    async def _parse_pptx(self, file_path: str) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(file_path)
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return ParsedDocument(
            text="\n".join(texts),
            metadata={"source": os.path.basename(file_path), "slide_count": len(prs.slides)},
            pages=None,
            tables=[],
        )

    async def _parse_csv(self, file_path: str) -> ParsedDocument:
        import csv

        rows: list[list[str]] = []
        text_parts: list[str] = []
        with open(file_path, newline="", encoding="utf-8-sig", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
                text_parts.append(" | ".join(row))
        return ParsedDocument(
            text="\n".join(text_parts),
            metadata={"source": os.path.basename(file_path), "row_count": len(rows)},
            pages=None,
            tables=[rows] if rows else [],
        )

    async def close(self) -> None:
        pass


async def _run_subprocess(cmd: list[str], capture: bool = False) -> str | bytes:
    proc = await asyncio.subprocess.create_subprocess_exec(*cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
    stdout, stderr = await proc.communicate()
    if proc.returncode != 0:
        raise RuntimeError(f"Command failed: {stderr.decode(errors='replace')}")
    if capture:
        return stdout.decode(errors="replace")
    return stdout


async def _read_text_file(file_path: str) -> str:
    """Read a text file in a thread."""
    return await asyncio.to_thread(lambda: open(file_path, "r", encoding="utf-8", errors="replace").read())


__all__ = ["UniversalDocumentParser", "DocumentParserFactory"]
